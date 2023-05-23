from pptx import Presentation
from pptx.util import Pt # This is used to set font size
import six
from google.cloud import translate_v2 as translate
import argostranslate.package
import argostranslate.translate
from translate import Translator
import ToAudio

DATABASE_DIR = "./database/"

def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text='AAA AAA'):
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

def translateToFile(path="", from_code = "ja", to_code = "en"):
    prs = Presentation(path)
    aryText = []
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            # If the shape doesn't have a text frame, continue to the next one
            if not shape.has_text_frame:
                continue
            # Loop through each paragraph in the shape's text frame
            for p, paragraph in enumerate(shape.text_frame.paragraphs):
                # Loop through each run in the paragraph
                for r, run in enumerate(paragraph.runs):
                    run.text = translate_text_2(run.text, from_code, to_code)
                    # run.font.size = Pt(14)  # Change the font size if needed
                    if run.text == "":
                        continue
                    aryText.append(run.text)
    prs.save(DATABASE_DIR + to_code + ".pptx")
    return aryText


def translate_text_online(text):
    # Replace 'YOUR_PROJECT_ID' with your actual project ID
    translate_client = translate.Client.from_service_account_json('path/to/your/credentials.json')
    result = translate_client.translate(text, source_language='ja', target_language='en')
    return result['translatedText']

def translate_text_1(text="hello", from_code = "ja", to_code = "en"):
    argostranslate.package.update_package_index()
    available_packages = argostranslate.package.get_available_packages()
    package_to_install = next(
        filter(
            lambda x: x.from_code == from_code and x.to_code == to_code, available_packages
        )
    )
    argostranslate.package.install_from_path(package_to_install.download())

    translation = argostranslate.translate.translate(text, from_code, to_code)
    return translation

def translate_text_2(text="hello", from_lang= "ja", to_lang= "en"):
    translator = Translator(to_lang=to_lang, from_lang=from_lang)
    if from_lang == "ja":
        text = text + "。"
    else:
        text = text + "."
    translation = translator.translate(text)
    translation = translation[0:len(translation) - 1]
    return translation

def saveToTxtFile(strSavePath, aryTxt):
    with open(strSavePath, "w", encoding="utf-8") as f:
        for txt in aryTxt:
            try:
                if len(txt) == 0:
                    continue
                f.write(txt + "\n")
            except:
                print('')
        f.close()

def main(strPath, strLang):
    aryTxt = translateToFile(path=strPath, from_code="ja", to_code=strLang)
    strSavePath = strLang + ".txt"
    saveToTxtFile(DATABASE_DIR + strSavePath, aryTxt)
    ToAudio.main(DATABASE_DIR + strSavePath, DATABASE_DIR + strLang + ".wav")
